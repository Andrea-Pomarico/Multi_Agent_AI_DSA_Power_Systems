"""
Agent 8 — Presentation Agent
==============================
Builds a PowerPoint (.pptx) presentation from the final report text and
simulation plot images.  Requires the `python-pptx` package; if it is not
installed the function returns (False, "", <error message>) gracefully so
the rest of the pipeline is unaffected.

Slide structure
---------------
  Slide 1        — Title slide (study-case name + generation timestamp)
  Slides 2..N    — Report text, 12 lines per slide
  Slide N+1      — Source Word report path (optional)
  Slides N+2..   — One image slide per PNG plot (full-width, 12.3 in)

Public API
----------
  presentation_agent(report_text, study_case, out_dir, label,
                     plot_paths=None, source_docx_path=None)
      → (ok: bool, pptx_path: str, message: str)
"""

import os
from datetime import datetime


def _chunk_lines(lines: list[str], chunk_size: int = 12) -> list[list[str]]:
    return [lines[i:i + chunk_size] for i in range(0, len(lines), chunk_size)]


def presentation_agent(report_text: str,
                       study_case: str,
                       out_dir: str,
                       label: str,
                       plot_paths: list[str] | None = None,
                       source_docx_path: str | None = None) -> tuple[bool, str, str]:
    """
    Build a PowerPoint presentation from final report text and plot images.
    Returns: (ok, pptx_path, message)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as e:
        return False, "", f"python-pptx not available: {e}"

    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, f"{label}_presentation.pptx")
    plot_paths = [p for p in (plot_paths or []) if os.path.isfile(p)]

    prs = Presentation()

    # Slide 1: title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"RMS Stability Report — {study_case}"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # Slide 2..n: report text chunks
    lines = [ln for ln in (report_text or "").splitlines() if ln.strip()]
    if not lines:
        lines = ["No report text available."]

    for idx, chunk in enumerate(_chunk_lines(lines, chunk_size=12), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Final Report (Part {idx})"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(chunk):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

    # Optional slide: source docx reference
    if source_docx_path:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Source Word Report"
        body = slide.shapes.placeholders[1].text_frame
        body.text = os.path.abspath(source_docx_path)
        body.paragraphs[0].font.size = Pt(14)

    # Image slides
    for img in plot_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = os.path.splitext(os.path.basename(img))[0].replace("_", " ").title()

        # Keep simple and robust placement for wide plots.
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(12.3)
        slide.shapes.add_picture(img, left, top, width=width)

    prs.save(pptx_path)
    return True, pptx_path, "Presentation created"
